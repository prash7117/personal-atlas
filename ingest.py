# Copyright (c) 2026 Prashanth Shankar Narayan
# SPDX-License-Identifier: Apache-2.0

import os
import re
import json
import time
import hashlib
from collections import Counter
import tempfile
from pathlib import Path
from typing import Dict, Iterable, List, Optional
from html.parser import HTMLParser

from dotenv import load_dotenv
from qdrant_client import QdrantClient
from qdrant_client.http.models import (
    Distance,
    PointStruct,
    PointIdsList,
    VectorParams,
)
from openai import OpenAI

# Optional PDF support:
#   pip install pypdf
try:
    from pypdf import PdfReader  # type: ignore
except Exception:
    PdfReader = None

# Optional DOCX support:
#   pip install python-docx
try:
    from docx import Document  # type: ignore
    from docx.oxml.ns import nsmap  # type: ignore
except Exception:
    Document = None
    nsmap = None

# Optional PPTX support:
#   pip install python-pptx
try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None

# Optional MSG support:
#   pip install extract-msg
try:
    import extract_msg  # type: ignore
except Exception:
    extract_msg = None

load_dotenv(dotenv_path=".env")

QDRANT_URL = os.environ["QDRANT_URL"]
COLLECTION = os.environ["QDRANT_COLLECTION"]
EMBED_MODEL = os.environ.get("OPENAI_EMBED_MODEL", "text-embedding-3-small")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")

if not OPENAI_API_KEY or "YOUR_KEY" in OPENAI_API_KEY:
    raise SystemExit(
        "OPENAI_API_KEY missing/placeholder. Set it in .env before running ingestion."
    )

# Tune these:
INCLUDE_EXT = {
    ".md",
    ".txt",
    ".rst",
    ".py",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".json",
    ".yaml",
    ".yml",
    ".sh",
    ".ini",
    ".cfg",
    ".pdf",  # works only if pypdf installed
    ".docx",  # works only if python-docx installed
    ".pptx",  # works only if python-pptx installed
    ".msg",  # works only if extract-msg installed
}

# Directories to skip anywhere in path
SKIP_DIRS = {
    ".git",
    ".hg",
    ".svn",
    ".venv",
    "venv",
    "__pycache__",
    "build",
    "dist",
    ".tox",
    ".mypy_cache",
    ".pytest_cache",
    ".idea",
    ".vscode",
}

# Files to skip by regex (edit as needed)
SKIP_FILE_RE = re.compile(r".*\.pyc$|.*\.pyo$|.*\.so$|.*\.a$|.*\.o$")

# Safety limits
MAX_FILE_BYTES = 2_000_000  # skip >2MB by default (raise if you want)
CHUNK_MAX_CHARS = 3500  # better for code than 2000
CHUNK_OVERLAP = 300
PDF_MARGIN_SCAN_LINES = 3
PDF_MARGIN_MIN_REPEAT_RATIO = 0.6

# MSG attachment extraction limits
MSG_ATTACHMENT_MAX_BYTES = 200_000
MSG_ATTACHMENT_MAX_COUNT = 25
MSG_ATTACHMENT_TEXT_EXT = {
    ".cfg",
    ".csv",
    ".ini",
    ".json",
    ".log",
    ".md",
    ".py",
    ".rst",
    ".sh",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}
MSG_ATTACHMENT_HTML_EXT = {".html", ".htm"}
MSG_ATTACHMENT_DOC_EXT = {".pdf", ".docx"}

# Batch sizes
EMBED_BATCH = 64  # embeddings input batch size
UPSERT_BATCH = 200  # qdrant upsert batch size

# Store ingestion state here (per collection/repo_id)
STATE_FILE = ".rag_ingest_state.json"
ROOT_ALIAS_PREFIX = "__root__"


def stable_point_id(repo_id: str, relpath: str, chunk_id: int) -> int:
    """
    Stable 64-bit integer ID based on (repo_id, relpath, chunk_id).
    This ensures upserts are deterministic across runs.
    """
    s = f"{repo_id}::{relpath}::{chunk_id}"
    h = hashlib.blake2b(s.encode("utf-8"), digest_size=8).digest()
    return int.from_bytes(h, "big", signed=False)


def sha1_text(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()


def _ensure_root_alias(
    root: str, primary_root: str, root_aliases: Dict[str, str]
) -> str:
    if root == primary_root:
        root_aliases[root] = ""
        return ""
    alias = root_aliases.get(root)
    if alias:
        return alias
    base = ROOT_ALIAS_PREFIX + sha1_text(root)[:8]
    alias = base
    existing = set(root_aliases.values())
    counter = 1
    while alias in existing or alias == "":
        alias = f"{base}_{counter}"
        counter += 1
    root_aliases[root] = alias
    return alias


def _relpath_key(
    root: str, relpath: str, primary_root: str, root_aliases: Dict[str, str]
) -> str:
    alias = _ensure_root_alias(root, primary_root, root_aliases)
    if alias:
        return alias + "/" + relpath
    return relpath


def sha1_file(path: Path) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def chunk_text(text: str, max_chars: int, overlap: int) -> List[str]:
    text = text.replace("\r\n", "\n")
    n = len(text)
    if n <= max_chars:
        return [text]

    chunks: List[str] = []
    i = 0
    while i < n:
        j = min(n, i + max_chars)
        chunks.append(text[i:j])
        if j == n:
            break
        i = max(0, j - overlap)
    return chunks


def _pdf_margin_signature(line: str) -> str:
    normalized = re.sub(r"\s+", " ", line.strip()).lower()
    if not normalized:
        return ""
    return re.sub(r"\d+", "#", normalized)


def _pdf_margin_signatures(lines: List[str], from_top: bool) -> List[str]:
    signatures: List[str] = []
    if from_top:
        indices = range(len(lines))
    else:
        indices = range(len(lines) - 1, -1, -1)

    for idx in indices:
        signature = _pdf_margin_signature(lines[idx])
        if not signature:
            continue
        signatures.append(signature)
        if len(signatures) >= PDF_MARGIN_SCAN_LINES:
            break
    return signatures


def _strip_pdf_headers_footers(page_texts: List[str]) -> List[str]:
    if len(page_texts) < 2:
        return page_texts

    pages_lines: List[List[str]] = []
    top_counts: Counter = Counter()
    bottom_counts: Counter = Counter()

    for page_text in page_texts:
        lines = page_text.replace("\r\n", "\n").split("\n")
        pages_lines.append(lines)
        top_counts.update(_pdf_margin_signatures(lines, from_top=True))
        bottom_counts.update(_pdf_margin_signatures(lines, from_top=False))

    min_repeat = max(2, int(len(page_texts) * PDF_MARGIN_MIN_REPEAT_RATIO))
    header_signatures = {
        signature for signature, count in top_counts.items() if count >= min_repeat
    }
    footer_signatures = {
        signature for signature, count in bottom_counts.items() if count >= min_repeat
    }
    if not header_signatures and not footer_signatures:
        return page_texts

    cleaned_pages: List[str] = []
    for lines in pages_lines:
        start = 0
        removed_top = 0
        while start < len(lines) and removed_top < PDF_MARGIN_SCAN_LINES:
            signature = _pdf_margin_signature(lines[start])
            if not signature:
                start += 1
                continue
            if signature in header_signatures:
                start += 1
                removed_top += 1
                continue
            break

        end = len(lines)
        removed_bottom = 0
        while end > start and removed_bottom < PDF_MARGIN_SCAN_LINES:
            signature = _pdf_margin_signature(lines[end - 1])
            if not signature:
                end -= 1
                continue
            if signature in footer_signatures:
                end -= 1
                removed_bottom += 1
                continue
            break

        cleaned_pages.append("\n".join(lines[start:end]))

    return cleaned_pages


def read_pdf(path: Path) -> str:
    if PdfReader is None:
        raise RuntimeError(
            "PDF support requested but pypdf is not installed. pip install pypdf"
        )
    reader = PdfReader(str(path))
    parts: List[str] = []
    for p in reader.pages:
        parts.append(p.extract_text() or "")
    cleaned_parts = _strip_pdf_headers_footers(parts)
    return "\n".join(cleaned_parts)


def read_docx(path: Path) -> str:
    if Document is None:
        raise RuntimeError(
            "DOCX support requested but python-docx is not installed. pip install python-docx"
        )
    doc = Document(str(path))
    parts: List[str] = []

    def append_paragraphs(paragraphs: Iterable) -> None:
        for p in paragraphs:
            if p.text:
                parts.append(p.text)

    def append_tables(tables: Iterable) -> None:
        for table in tables:
            for row in table.rows:
                row_text: List[str] = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    parts.append(" | ".join(row_text))

    def append_textboxes_from_element(element: object) -> None:
        if nsmap is None:
            return
        try:
            tbx_paras = element.xpath(".//w:txbxContent//w:p", namespaces=nsmap)
        except Exception:
            return
        for p in tbx_paras:
            texts: List[str] = []
            try:
                text_nodes = p.xpath(".//w:t", namespaces=nsmap)
            except Exception:
                text_nodes = []
            for node in text_nodes:
                if node.text:
                    texts.append(node.text)
            if texts:
                parts.append("".join(texts))

    append_paragraphs(doc.paragraphs)
    append_tables(doc.tables)

    for section in doc.sections:
        append_paragraphs(section.header.paragraphs)
        append_tables(section.header.tables)
        append_paragraphs(section.footer.paragraphs)
        append_tables(section.footer.tables)

    doc_element = getattr(doc, "_element", None)
    if doc_element is None:
        doc_element = getattr(doc, "element", None)
    if doc_element is not None:
        append_textboxes_from_element(doc_element)

    for section in doc.sections:
        for header_footer in (section.header, section.footer):
            element = getattr(header_footer, "_element", None)
            if element is None:
                element = getattr(header_footer, "element", None)
            if element is not None:
                append_textboxes_from_element(element)

    return "\n".join(parts)


def read_pptx(path: Path) -> str:
    if Presentation is None:
        raise RuntimeError(
            "PPTX support requested but python-pptx is not installed. pip install python-pptx"
        )
    pptx = Presentation(str(path))
    parts: List[str] = []
    slide_num = 0
    for slide in pptx.slides:
        slide_num += 1
        parts.append(f"Slide {slide_num}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = getattr(shape, "text", "")
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    row_text: List[str] = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        parts.append(" | ".join(row_text))
        parts.append("")
    return "\n".join(parts).strip()


def read_msg(path: Path) -> str:
    if extract_msg is None:
        raise RuntimeError(
            "MSG support requested but extract-msg is not installed. pip install extract-msg"
        )
    msg = extract_msg.Message(str(path))

    class _HTMLTextParser(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: List[str] = []
            self._last_newline = False

        def handle_starttag(self, tag: str, attrs) -> None:
            if tag in ("br", "p", "div", "li", "tr", "td", "th"):
                self._append_newline()

        def handle_endtag(self, tag: str) -> None:
            if tag in ("p", "div", "li", "tr", "table", "ul", "ol"):
                self._append_newline()

        def handle_data(self, data: str) -> None:
            if not data:
                return
            self.parts.append(data)
            self._last_newline = False

        def error(self, message: str) -> None:
            _ = message

        def _append_newline(self) -> None:
            if not self.parts:
                self.parts.append("\n")
                self._last_newline = True
                return
            if not self._last_newline:
                self.parts.append("\n")
                self._last_newline = True

    def _norm(value: object) -> str:
        if value is None:
            return ""
        if isinstance(value, bytes):
            return value.decode(errors="ignore")
        return str(value)

    def _html_to_text(value: str) -> str:
        parser = _HTMLTextParser()
        parser.feed(value)
        parser.close()
        text = "".join(parser.parts)
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        while "\n\n\n" in text:
            text = text.replace("\n\n\n", "\n\n")
        return text.strip()

    def _read_attachment_file(data: bytes, suffix: str) -> str:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(data)
                tmp_path = tmp.name
            if suffix == ".pdf":
                return read_pdf(Path(tmp_path))
            if suffix == ".docx":
                return read_docx(Path(tmp_path))
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
        return ""

    def _attachment_bytes(attachment: object) -> Optional[bytes]:
        data = getattr(attachment, "data", None)
        if data is None and hasattr(attachment, "get_data"):
            try:
                data = attachment.get_data()
            except Exception:
                data = None
        if data is None:
            return None
        if isinstance(data, bytes):
            return data
        if isinstance(data, bytearray):
            return bytes(data)
        return None

    try:
        parts: List[str] = []
        subject = _norm(getattr(msg, "subject", None))
        sender = _norm(getattr(msg, "sender", None))
        to = _norm(getattr(msg, "to", None))
        cc = _norm(getattr(msg, "cc", None))
        date = _norm(getattr(msg, "date", None))

        if subject:
            parts.append("Subject: " + subject)
        if sender:
            parts.append("From: " + sender)
        if to:
            parts.append("To: " + to)
        if cc:
            parts.append("Cc: " + cc)
        if date:
            parts.append("Date: " + date)

        attachments: List[str] = []
        attachment_texts: List[str] = []
        attachment_count = 0
        for attachment in getattr(msg, "attachments", []) or []:
            name = (
                _norm(getattr(attachment, "longFilename", None))
                or _norm(getattr(attachment, "shortFilename", None))
                or _norm(getattr(attachment, "filename", None))
            )
            if name:
                attachments.append(name)
            if attachment_count >= MSG_ATTACHMENT_MAX_COUNT:
                continue
            if name:
                _, ext = os.path.splitext(name)
                ext = ext.lower()
                if ext in MSG_ATTACHMENT_TEXT_EXT or ext in MSG_ATTACHMENT_HTML_EXT:
                    data = _attachment_bytes(attachment)
                    if data is not None and len(data) <= MSG_ATTACHMENT_MAX_BYTES:
                        content = data.decode(errors="ignore")
                        if ext in MSG_ATTACHMENT_HTML_EXT:
                            content = _html_to_text(content)
                        attachment_texts.append("Attachment: " + name + "\n" + content)
                        attachment_count += 1
                elif ext in MSG_ATTACHMENT_DOC_EXT:
                    data = _attachment_bytes(attachment)
                    if data is not None and len(data) <= MSG_ATTACHMENT_MAX_BYTES:
                        try:
                            content = _read_attachment_file(data, ext)
                        except Exception:
                            content = ""
                        if content:
                            attachment_texts.append(
                                "Attachment: " + name + "\n" + content
                            )
                            attachment_count += 1
        if attachments:
            parts.append("Attachments: " + ", ".join(attachments))

        body = _norm(getattr(msg, "body", None))
        if body:
            parts.append("")
            parts.append(body)
        else:
            html_body = _norm(
                getattr(msg, "htmlBody", None) or getattr(msg, "html_body", None)
            )
            if html_body:
                parts.append("")
                parts.append(html_body)

        if attachment_texts:
            parts.append("")
            parts.append("Attachment contents:")
            parts.extend(attachment_texts)

        return "\n".join(parts).strip()
    finally:
        try:
            msg.close()
        except Exception:
            pass


def read_file_text(path: Path) -> str:
    suf = path.suffix.lower()
    if suf == ".pdf":
        return read_pdf(path)
    if suf == ".docx":
        return read_docx(path)
    if suf == ".pptx":
        return read_pptx(path)
    if suf == ".msg":
        return read_msg(path)
    return path.read_text(errors="ignore")


def should_skip_path(path: Path) -> bool:
    # Skip if any parent directory is in SKIP_DIRS
    for part in path.parts:
        if part in SKIP_DIRS:
            return True
    if SKIP_FILE_RE.match(path.name):
        return True
    return False


def load_state() -> Dict[str, Dict]:
    if not Path(STATE_FILE).exists():
        return {}
    try:
        return json.loads(Path(STATE_FILE).read_text())
    except Exception:
        return {}


def save_state(state: Dict[str, Dict]) -> None:
    Path(STATE_FILE).write_text(json.dumps(state, indent=2, sort_keys=True))


def is_indexable_file(p: Path) -> bool:
    if not p.is_file():
        return False
    if should_skip_path(p):
        return False
    if p.suffix.lower() not in INCLUDE_EXT:
        return False
    try:
        st = p.stat()
    except Exception:
        return False
    if st.st_size > MAX_FILE_BYTES:
        return False
    return True


def iter_indexable_files(
    base_root: Path, only_file: Optional[Path] = None
) -> Iterable[Path]:
    """
    If only_file is provided, yields that file (if indexable).
    Otherwise recursively yields indexable files under base_root.
    """
    if only_file is not None:
        if is_indexable_file(only_file):
            yield only_file
        return

    for p in base_root.rglob("*"):
        if is_indexable_file(p):
            yield p


def embed_texts(oai: OpenAI, texts: List[str]) -> List[List[float]]:
    vectors: List[List[float]] = []
    for i in range(0, len(texts), EMBED_BATCH):
        batch = texts[i : i + EMBED_BATCH]
        resp = oai.embeddings.create(model=EMBED_MODEL, input=batch)
        vectors.extend([d.embedding for d in resp.data])
    return vectors


def ensure_collection_exists(qdrant: QdrantClient, oai: OpenAI) -> None:
    """Create collection lazily when missing, using current embed model dimension."""
    if qdrant.collection_exists(COLLECTION):
        return

    probe = oai.embeddings.create(model=EMBED_MODEL, input=["qdrant-probe"])
    vector_size = len(probe.data[0].embedding)

    try:
        qdrant.create_collection(
            collection_name=COLLECTION,
            vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
        )
        print(
            f"Created missing collection '{COLLECTION}' "
            f"(vector_size={vector_size})."
        )
    except Exception as exc:
        # Another process may create it concurrently between the exists check and create call.
        if not qdrant.collection_exists(COLLECTION):
            raise exc


def upsert_points(qdrant: QdrantClient, points: List[PointStruct]) -> None:
    if points:
        qdrant.upsert(collection_name=COLLECTION, points=points)


def delete_points_for_relpaths(
    qdrant: QdrantClient,
    repo_id: str,
    relpaths: List[str],
    prev_files: Dict[str, Dict],
    batch: int = 200,
) -> None:
    """Delete all points for relpaths by computing stable point IDs."""
    if not relpaths:
        return

    ids: List[int] = []
    deleted_files = 0
    for relpath in relpaths:
        meta = prev_files.get(relpath)
        num_chunks = 0
        if isinstance(meta, dict):
            num_chunks = int(meta.get("num_chunks", 0) or 0)
        if num_chunks <= 0:
            deleted_files += 1
            continue
        for chunk_id in range(num_chunks):
            ids.append(stable_point_id(repo_id, relpath, chunk_id))
            if len(ids) >= batch:
                qdrant.delete(
                    collection_name=COLLECTION,
                    points_selector=PointIdsList(points=ids),
                )
                ids = []
        deleted_files += 1

    if ids:
        qdrant.delete(
            collection_name=COLLECTION,
            points_selector=PointIdsList(points=ids),
        )

    print(f"Deleted vectors for {deleted_files} missing files")


def delete_stale_chunks_for_file(
    qdrant: QdrantClient,
    repo_id: str,
    relpath: str,
    from_chunk_id: int,
    to_chunk_id_inclusive: int,
) -> None:
    """
    Delete stale chunk vectors for a single file by computing stable point IDs.
    Deletes chunk_id range [from_chunk_id, to_chunk_id_inclusive].
    """
    if to_chunk_id_inclusive < from_chunk_id:
        return

    ids = [
        stable_point_id(repo_id, relpath, cid)
        for cid in range(from_chunk_id, to_chunk_id_inclusive + 1)
    ]

    qdrant.delete(
        collection_name=COLLECTION,
        points_selector=PointIdsList(points=ids),
    )


def ingest(
    path_str: str, repo_id: Optional[str] = None, delete_missing: bool = False
) -> None:
    target = Path(path_str).resolve()
    if not target.exists():
        raise SystemExit(f"Path does not exist: {target}")

    # Support both directory and single file
    if target.is_file():
        base_root = target.parent
        only_file = target
        # Safety: do not delete missing based on a single-file scan
        if delete_missing:
            print("Note: --delete-missing ignored for single-file ingestion.")
            delete_missing = False
    else:
        base_root = target
        only_file = None

    if repo_id is None:
        # stable-ish id derived from base_root path
        repo_id = sha1_text(str(base_root))[:10]

    state = load_state()
    key = f"{COLLECTION}::{repo_id}"
    prev = state.get(key, {"files": {}})
    prev_files: Dict[str, Dict] = prev.get("files", {})

    base_root_str = str(base_root)
    primary_root = None
    roots: List[str] = []
    root_aliases: Dict[str, str] = {}
    if isinstance(prev, dict):
        primary_root = prev.get("root")
        roots_value = prev.get("roots")
        if isinstance(roots_value, list):
            roots = roots_value
        aliases_value = prev.get("root_aliases")
        if isinstance(aliases_value, dict):
            root_aliases = aliases_value
    if not primary_root:
        primary_root = base_root_str
    roots = [r for r in roots if isinstance(r, str)]
    if primary_root not in roots:
        roots.insert(0, primary_root)
    if base_root_str not in roots:
        roots.append(base_root_str)
    if not isinstance(root_aliases, dict):
        root_aliases = {}
    root_aliases[primary_root] = ""
    for root in roots:
        if root != primary_root:
            _ensure_root_alias(root, primary_root, root_aliases)

    alias_prefix_map: Dict[str, str] = {}
    for root, alias in root_aliases.items():
        if alias:
            alias_prefix_map[alias + "/"] = root

    qdrant = QdrantClient(url=QDRANT_URL)
    oai = OpenAI(api_key=OPENAI_API_KEY)
    ensure_collection_exists(qdrant, oai)

    seen_relpaths = set()
    points_buffer: List[PointStruct] = []

    changed_files = 0
    total_chunks = 0
    skipped_unchanged = 0

    started = time.time()

    for path in iter_indexable_files(base_root, only_file=only_file):
        relpath_display = str(path.relative_to(base_root))
        relpath_key = _relpath_key(
            base_root_str, relpath_display, primary_root, root_aliases
        )
        seen_relpaths.add(relpath_key)

        try:
            mtime = int(path.stat().st_mtime)
        except Exception:
            continue

        prev_meta = prev_files.get(relpath_key)

        # Fast skip if mtime unchanged
        if prev_meta and prev_meta.get("mtime") == mtime:
            if isinstance(prev_meta, dict):
                needs_root = prev_meta.get("root") != base_root_str
                needs_relpath = prev_meta.get("relpath") != relpath_display
                if needs_root or needs_relpath:
                    updated_meta = dict(prev_meta)
                    updated_meta["root"] = base_root_str
                    updated_meta["relpath"] = relpath_display
                    prev_files[relpath_key] = updated_meta
            skipped_unchanged += 1
            continue

        # Hash file content for incremental correctness
        try:
            content_hash = sha1_file(path)
        except Exception:
            continue

        # If content unchanged, update mtime and skip
        if prev_meta and prev_meta.get("sha1") == content_hash:
            # Preserve prior num_chunks if present (needed for stale chunk cleanup)
            num_chunks_prev = 0
            if isinstance(prev_meta, dict):
                num_chunks_prev = int(prev_meta.get("num_chunks", 0) or 0)

            prev_files[relpath_key] = {
                "mtime": mtime,
                "sha1": content_hash,
                "num_chunks": num_chunks_prev,
                "root": base_root_str,
                "relpath": relpath_display,
            }
            skipped_unchanged += 1
            continue

        # Read + chunk + embed
        try:
            text = read_file_text(path)
        except Exception as e:
            print(f"Skip (read failed): {relpath_display} -> {e}")
            continue

        if not text.strip():
            # If file becomes empty, delete all previously-indexed chunks for this file
            old_n = 0
            if isinstance(prev_meta, dict):
                old_n = int(prev_meta.get("num_chunks", 0) or 0)

            if old_n > 0:
                delete_stale_chunks_for_file(qdrant, repo_id, relpath_key, 0, old_n - 1)
                print(
                    "Deleted stale chunks for empty file: "
                    f"{relpath_display} (0..{old_n - 1})"
                )

            prev_files[relpath_key] = {
                "mtime": mtime,
                "sha1": content_hash,
                "num_chunks": 0,
                "root": base_root_str,
                "relpath": relpath_display,
            }
            changed_files += 1
            continue

        chunks = chunk_text(text, CHUNK_MAX_CHARS, CHUNK_OVERLAP)
        new_n = len(chunks)
        vectors = embed_texts(oai, chunks)

        for idx, (chunk, vec) in enumerate(zip(chunks, vectors)):
            pid = stable_point_id(repo_id, relpath_key, idx)
            points_buffer.append(
                PointStruct(
                    id=pid,
                    vector=vec,
                    payload={
                        "repo_id": repo_id,
                        "root": base_root_str,
                        "relpath": relpath_display,
                        "path": str(path),
                        "chunk_id": idx,
                        "ext": path.suffix.lower(),
                        "mtime": mtime,
                        "sha1": content_hash,
                        "text": chunk,
                    },
                )
            )
            total_chunks += 1

        # If the file shrank, delete stale tail chunks from the previous version
        old_n = 0
        if isinstance(prev_meta, dict):
            old_n = int(prev_meta.get("num_chunks", 0) or 0)

        if old_n > new_n:
            delete_stale_chunks_for_file(qdrant, repo_id, relpath_key, new_n, old_n - 1)
            print(f"Deleted stale chunks for {relpath_display}: {new_n}..{old_n - 1}")

        prev_files[relpath_key] = {
            "mtime": mtime,
            "sha1": content_hash,
            "num_chunks": new_n,
            "root": base_root_str,
            "relpath": relpath_display,
        }
        changed_files += 1

        if len(points_buffer) >= UPSERT_BATCH:
            upsert_points(qdrant, points_buffer)
            points_buffer.clear()
            print(
                f"Upserted... changed_files={changed_files}, chunks={total_chunks}, "
                f"skipped_unchanged={skipped_unchanged}"
            )

    # final flush
    if points_buffer:
        upsert_points(qdrant, points_buffer)

    # Delete missing files (directory ingestion only)
    if delete_missing:
        missing: List[str] = []
        for relpath_key, meta in prev_files.items():
            if relpath_key in seen_relpaths:
                continue
            meta_root = None
            if isinstance(meta, dict):
                meta_root = meta.get("root")
            if not meta_root:
                for prefix, root in alias_prefix_map.items():
                    if relpath_key.startswith(prefix):
                        meta_root = root
                        break
            if not meta_root:
                meta_root = primary_root
            if meta_root != base_root_str:
                continue
            missing.append(relpath_key)
        if missing:
            print(
                f"delete_missing: {len(missing)} files missing; deleting vectors from Qdrant..."
            )
            delete_points_for_relpaths(qdrant, repo_id, missing, prev_files, batch=200)
            for rp in missing:
                prev_files.pop(rp, None)

    display_name = None
    if isinstance(prev, dict) and "display_name" in prev:
        display_name = prev.get("display_name")

    entry = {
        "files": prev_files,
        "root": primary_root,
        "roots": roots,
        "root_aliases": root_aliases,
        "repo_id": repo_id,
        "last_run_ts": time.time(),
    }
    if display_name is not None:
        entry["display_name"] = display_name
    state[key] = entry
    save_state(state)

    elapsed = time.time() - started
    print(
        f"DONE. changed_files={changed_files}, chunks={total_chunks}, "
        f"skipped_unchanged={skipped_unchanged}, seconds={elapsed:.1f}"
    )


if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser()
    ap.add_argument("path", help="directory OR single file to index")
    ap.add_argument("--repo-id", default=None, help="stable repo id label (optional)")
    ap.add_argument(
        "--delete-missing",
        action="store_true",
        help="for directory ingestion: delete vectors for files that disappeared since last full scan",
    )
    args = ap.parse_args()

    ingest(args.path, repo_id=args.repo_id, delete_missing=args.delete_missing)
